from pptx import Presentation

def extrair_texto_pptx(caminho_pptx: str):
    """Extrai o texto de um arquivo PPTX e separa por parágrafos."""

    apresentacao = Presentation(caminho_pptx)
    informacoes = []

    numero_paragrafo = 1

    for numero_slide, slide in enumerate(apresentacao.slides, start=1):

        for shape in slide.shapes:

            if not shape.has_text_frame:
                continue

            for paragrafo in shape.text_frame.paragraphs:

                texto = paragrafo.text.strip()

                if not texto:
                    continue

                informacoes.append({
                    "texto": texto,
                    "metadados": {
                        "fonte": caminho_pptx,
                        "paragrafo": numero_paragrafo,
                        "slide": numero_slide
                    }
                })

                numero_paragrafo += 1

    return informacoes